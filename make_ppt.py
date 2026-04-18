"""
Generate FinSight_Presentation.pptx  (9 slides)
Run:  python make_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour palette ─────────────────────────────────────────────────────────────
NAVY   = RGBColor(0x0d, 0x1b, 0x4b)
BLUE   = RGBColor(0x1c, 0x3d, 0x7a)
INDIGO = RGBColor(0x4c, 0x6e, 0xf5)
LTBLUE = RGBColor(0xdd, 0xe9, 0xff)
LTGREY = RGBColor(0xf4, 0xf6, 0xfb)
WHITE  = RGBColor(0xff, 0xff, 0xff)
GREY   = RGBColor(0x88, 0x92, 0xaa)
GREEN  = RGBColor(0x2d, 0x9c, 0x5a)
ORANGE = RGBColor(0xe6, 0x8a, 0x00)
RED    = RGBColor(0xc0, 0x39, 0x2b)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ── Helpers ────────────────────────────────────────────────────────────────────
def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=None):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.background() if not fill else (shape.fill.solid() or setattr(shape.fill.fore_color, 'rgb', fill))
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
        if line_w:
            shape.line.width = Pt(line_w)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h, size=14, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb

def add_para(tf, text, size=11, bold=False, color=NAVY,
             align=PP_ALIGN.LEFT, italic=False, space_before=0):
    p = tf.add_paragraph()
    p.alignment = align
    p.space_before = Pt(space_before)
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return p

def header(slide, title, subtitle=None):
    add_rect(slide, 0, 0, 13.33, 1.3, fill=NAVY)
    add_rect(slide, 0, 1.3, 13.33, 0.05, fill=INDIGO)
    add_text(slide, title, 0.45, 0.1, 12.5, 0.7, size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle, 0.45, 0.78, 12.5, 0.45, size=13, color=LTBLUE)

def foot(slide, note="FinSight — Cloud Computing Project"):
    add_rect(slide, 0, 7.22, 13.33, 0.28, fill=NAVY)
    add_text(slide, note, 0.4, 7.24, 12.5, 0.22, size=8, color=GREY)

def bg(slide, color=WHITE):
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide, NAVY)
add_rect(slide, 0, 0, 0.2, 7.5, fill=INDIGO)

add_text(slide, "FinSight", 0.6, 1.3, 12, 1.3,
         size=72, bold=True, color=WHITE)
add_text(slide, "Financial Document Intelligence Platform",
         0.6, 2.8, 11, 0.65, size=24, color=LTBLUE)
add_rect(slide, 0.6, 3.6, 4.0, 0.06, fill=INDIGO)
add_text(slide, "Upload receipts, invoices & bank statements.\nAI extracts, categorises and visualises your finances — automatically.",
         0.6, 3.75, 10, 0.9, size=15, color=GREY)
add_text(slide, "Cloud Computing Project  ·  AWS Academy",
         0.6, 5.0, 8, 0.45, size=13, italic=True, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem & Solution
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header(slide, "The Problem & Our Solution")
foot(slide)

# Problem
add_rect(slide, 0.4, 1.55, 5.9, 5.55, fill=LTGREY, line=RGBColor(0xc5,0xd3,0xf0), line_w=0.6)
add_rect(slide, 0.4, 1.55, 5.9, 0.5, fill=RED)
add_text(slide, "  The Problem", 0.45, 1.58, 5.5, 0.4, size=15, bold=True, color=WHITE)

tb = slide.shapes.add_textbox(Inches(0.6), Inches(2.2), Inches(5.5), Inches(4.7))
tb.word_wrap = True
tf = tb.text_frame
tf.word_wrap = True
problems = [
    "People receive financial documents in many formats — PDFs, scanned images, photos of receipts.",
    "Manually reading, typing out amounts, dates and vendor names is slow and error-prone.",
    "Tracking monthly spend and income across tens of documents is nearly impossible without a system.",
    "Existing tools like Excel require manual entry — there is no automatic extraction.",
]
first = True
for pr in problems:
    if first:
        p = tf.paragraphs[0]; first = False
    else:
        p = tf.add_paragraph()
        p.space_before = Pt(8)
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "✗  " + pr
    r.font.size = Pt(11)
    r.font.color.rgb = NAVY

# Solution
add_rect(slide, 6.7, 1.55, 6.2, 5.55, fill=LTGREY, line=RGBColor(0xc5,0xd3,0xf0), line_w=0.6)
add_rect(slide, 6.7, 1.55, 6.2, 0.5, fill=GREEN)
add_text(slide, "  Our Solution — FinSight", 6.75, 1.58, 5.8, 0.4, size=15, bold=True, color=WHITE)

tb2 = slide.shapes.add_textbox(Inches(6.9), Inches(2.2), Inches(5.8), Inches(4.7))
tb2.word_wrap = True
tf2 = tb2.text_frame
tf2.word_wrap = True
solutions = [
    ("Upload any document", "PDF, PNG, JPG — up to 50 MB per file, multiple files at once."),
    ("Automatic OCR", "Tesseract reads the text from scanned pages, even handwritten receipts."),
    ("Smart extraction", "Pulls out amounts, dates, vendor names, invoice numbers automatically."),
    ("Auto-categorisation", "Classifies each document: Invoice, Receipt, Pay Stub, Bank Statement, Tax Form."),
    ("Live dashboard", "Monthly cash flow chart, top vendors, category breakdown — all filtered by date."),
]
first = True
for title, desc in solutions:
    if first:
        p = tf2.paragraphs[0]; first = False
    else:
        p = tf2.add_paragraph()
        p.space_before = Pt(6)
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = "✓  " + title
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = GREEN

    p2 = tf2.add_paragraph()
    p2.alignment = PP_ALIGN.LEFT
    r2 = p2.add_run()
    r2.text = "     " + desc
    r2.font.size = Pt(10)
    r2.font.color.rgb = NAVY


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Product Walkthrough
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header(slide, "Product Walkthrough", "Step by step — what the user sees and does")
foot(slide)

steps = [
    ("1", "Register / Login",
     "Create an account with username + password.\nAll documents are private per user.\nPasswords hashed with Werkzeug.",
     LTBLUE, INDIGO),
    ("2", "Upload Documents",
     "Pick one or multiple files at once.\nProgress bar shows per-file status.\nSupports PDF, PNG, JPG, TIFF.",
     LTBLUE, INDIGO),
    ("3", "Automatic OCR",
     "Tesseract reads text at 300 DPI.\nExtracts: amounts, dates, vendor,\ninvoice number, document type.",
     LTBLUE, INDIGO),
    ("4", "Review Documents",
     "See all uploaded documents with status.\nClick any to view extracted fields\nand preview original file.",
     LTBLUE, INDIGO),
    ("5", "Dashboard",
     "Income vs Expenses cards.\nMonthly cash flow trend chart.\nTop vendors bar chart.\nCategory breakdown.",
     LTBLUE, INDIGO),
    ("6", "Export / Filter",
     "Filter dashboard by month (e.g. Jan 2026).\nExport all documents as CSV.\nReprocess any document with one click.",
     LTBLUE, INDIGO),
]

box_w = 3.95
box_h = 4.5
positions = [(0.35, 1.55), (4.35, 1.55), (8.35, 1.55)]

# top row
for idx in range(3):
    sx, sy = positions[idx]
    num, title, desc, fill_c, border_c = steps[idx]
    add_rect(slide, sx, sy, box_w, box_h, fill=fill_c, line=border_c, line_w=0.8)
    # number circle area
    add_rect(slide, sx, sy, box_w, 0.52, fill=NAVY)
    add_text(slide, f"Step {num}  —  {title}", sx + 0.15, sy + 0.08, box_w - 0.25, 0.38,
             size=13, bold=True, color=WHITE)
    add_text(slide, desc, sx + 0.2, sy + 0.68, box_w - 0.4, 3.6,
             size=11, color=NAVY)

# bottom row — show as horizontal strip
add_rect(slide, 0.35, 6.2, 12.6, 0.85, fill=NAVY)
add_text(slide, "What makes FinSight different", 0.5, 6.22, 4.0, 0.38, size=12, bold=True, color=WHITE)
differentiators = [
    "No manual data entry",
    "Works on scanned images, not just digital PDFs",
    "Multi-user — each person sees only their documents",
    "Runs entirely on AWS — accessible from any browser",
]
xx = 3.8
for d in differentiators:
    add_text(slide, "✓  " + d, xx, 6.25, 2.5, 0.75, size=10, color=LTBLUE)
    xx += 2.55


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — How OCR & Extraction Works
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header(slide, "How OCR & Data Extraction Works", "The intelligence behind FinSight")
foot(slide)

# pipeline flow
add_rect(slide, 0.35, 1.55, 12.65, 0.42, fill=NAVY)
add_text(slide, "Document Processing Pipeline", 0.5, 1.58, 12.0, 0.32, size=12, bold=True, color=WHITE)

pipe_steps = [
    ("PDF / Image\nUploaded", "File received\nby Flask API"),
    ("pdf2image\nconverts PDF", "Each page →\nPNG at 300 DPI"),
    ("Tesseract OCR\nreads image", "Outputs raw\ntext string"),
    ("Normalise\ntext", "Fix OCR\nartefacts\ne.g. 'J anuary'"),
    ("Regex\nextraction", "Amounts, dates,\nvendor, invoice\nnumber"),
    ("Classify\ndocument", "Invoice / Receipt\nPay Stub / Bank\nStatement / Tax"),
    ("Save to\nPostgreSQL", "All fields stored\nwith confidence\nscore"),
]
bw = 1.73
xx = 0.35
for i, (top, bot) in enumerate(pipe_steps):
    add_rect(slide, xx, 2.05, bw - 0.08, 1.55, fill=LTBLUE, line=INDIGO, line_w=0.8)
    add_text(slide, top, xx + 0.1, 2.1, bw - 0.25, 0.55, size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_rect(slide, xx, 2.6, bw - 0.08, 0.04, fill=INDIGO)
    add_text(slide, bot, xx + 0.1, 2.68, bw - 0.25, 0.85, size=9, color=NAVY, align=PP_ALIGN.CENTER)
    if i < len(pipe_steps) - 1:
        add_text(slide, "→", xx + bw - 0.08, 2.65, 0.15, 0.4, size=13, bold=True, color=INDIGO)
    xx += bw

# extraction details
add_rect(slide, 0.35, 3.82, 12.65, 0.42, fill=NAVY)
add_text(slide, "Fields Extracted from Every Document", 0.5, 3.85, 12.0, 0.32, size=12, bold=True, color=WHITE)

fields = [
    ("Amounts", "$1,234.56\n$0.99 / 12,000 USD", "Regex: all dollar\nvalues found"),
    ("Total Amount", "Largest / labelled\ntotal in document", "4-stage logic:\nlabel → sum → max"),
    ("Dates", "01/15/2026\nJanuary 15, 2026", "Matches 6 different\ndate formats"),
    ("Vendor Name", "First line or line\nafter 'From:' / 'Vendor:'", "Signal-word\nlookup + fallback"),
    ("Invoice Number", "INV-0042\nOrder #8821", "Matches invoice,\nreceipt, ref, txn"),
    ("Document Type", "Invoice / Receipt\nPay Stub / Bank Stmt", "Keyword scoring\nacross 6 categories"),
    ("Confidence", "0.0 – 1.0 score", "Ratio of readable\nchars to total chars"),
]
bw2 = 1.73
xx = 0.35
for fname, example, how in fields:
    add_rect(slide, xx, 4.35, bw2 - 0.06, 2.65, fill=LTGREY, line=RGBColor(0xc5,0xd3,0xf0), line_w=0.5)
    add_text(slide, fname,   xx+0.1, 4.4,  bw2-0.2, 0.38, size=10, bold=True, color=NAVY)
    add_text(slide, example, xx+0.1, 4.82, bw2-0.2, 0.75, size=9,  color=INDIGO, italic=True)
    add_rect(slide, xx, 5.58, bw2-0.06, 0.03, fill=LTBLUE)
    add_text(slide, how,     xx+0.1, 5.65, bw2-0.2, 1.2, size=8, color=GREY)
    xx += bw2


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Dashboard & Features
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header(slide, "Dashboard & Key Features", "What users see after uploading documents")
foot(slide)

# stat cards row
cards = [
    ("Total Income",    "$8,450",  "Sum of all income docs",   GREEN),
    ("Total Expenses",  "$3,210",  "Sum of all expense docs",  RED),
    ("Net Cash Flow",   "$5,240",  "Income minus expenses",    INDIGO),
    ("Documents",       "24",      "Total uploaded this month",BLUE),
    ("Savings Rate",    "62%",     "Net / Income ratio",       ORANGE),
    ("Avg Monthly",     "$1,605",  "Average monthly spend",    NAVY),
]
cw = 2.05
xx = 0.35
for label, val, desc, color in cards:
    add_rect(slide, xx, 1.55, cw - 0.08, 1.35, fill=LTGREY, line=color, line_w=1.2)
    add_rect(slide, xx, 1.55, cw - 0.08, 0.08, fill=color)
    add_text(slide, val,   xx+0.12, 1.72, cw-0.25, 0.55, size=22, bold=True, color=color)
    add_text(slide, label, xx+0.12, 2.28, cw-0.25, 0.32, size=10, bold=True, color=NAVY)
    add_text(slide, desc,  xx+0.12, 2.6,  cw-0.25, 0.25, size=8,  color=GREY)
    xx += cw

# three feature panels
panels = [
    ("Monthly Cash Flow Trend",
     "Bar chart showing income and expenses for each of the last 12 months.\n\n"
     "Groups documents by the date extracted from the document itself (primary_date) — not the upload date.\n\n"
     "This means uploading January documents in April still shows them correctly in January.",
     "Chart type: grouped bar\nX-axis: YYYY-MM\nY-axis: total amount ($)"),
    ("Top Vendors by Spend",
     "Horizontal bar chart ranking vendors by total amount spent.\n\n"
     "Vendor name is extracted by OCR — looks for signal words like 'From:', 'Vendor:', 'Billed by:' and falls back to the first line of the document.\n\n"
     "Sorted by SUM of total_amount, not by document count.",
     "Chart type: horizontal bar\nMetric: total spend ($)\nTop 8 vendors shown"),
    ("Category Breakdown",
     "Donut/pie breakdown of documents by type: Invoice, Receipt, Pay Stub, Bank Statement, Tax Form.\n\n"
     "Each document is classified by a keyword scoring system — the category with the most keyword matches wins.\n\n"
     "Filterable by month using the month picker.",
     "6 document categories\nMonth filter: YYYY-MM\nCSV export available"),
]
pw = 4.18
xx = 0.35
for ptitle, pdesc, pnote in panels:
    add_rect(slide, xx, 3.12, pw - 0.08, 3.95, fill=LTGREY, line=RGBColor(0xc5,0xd3,0xf0), line_w=0.6)
    add_rect(slide, xx, 3.12, pw - 0.08, 0.45, fill=NAVY)
    add_text(slide, ptitle, xx+0.12, 3.15, pw-0.25, 0.36, size=12, bold=True, color=WHITE)
    add_text(slide, pdesc,  xx+0.12, 3.68, pw-0.25, 2.4,  size=10, color=NAVY)
    add_rect(slide, xx, 6.1, pw-0.08, 0.04, fill=INDIGO)
    add_text(slide, pnote,  xx+0.12, 6.18, pw-0.25, 0.75, size=8, italic=True, color=GREY)
    xx += pw


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header(slide, "System Architecture", "Three Docker containers on one EC2, backed by S3")
foot(slide)

# Browser
add_rect(slide, 0.3, 1.65, 2.1, 1.1, fill=LTBLUE, line=INDIGO, line_w=1.0)
add_text(slide, "🌐  Browser", 0.45, 1.73, 1.9, 0.38, size=12, bold=True, color=NAVY)
add_text(slide, "Any device\nHTTP port 80", 0.45, 2.1, 1.9, 0.5, size=9, color=BLUE)

add_text(slide, "HTTP :80  →", 2.45, 2.05, 1.35, 0.35, size=9, italic=True, color=GREY)

# EC2 outer box
add_rect(slide, 3.65, 1.5, 6.4, 5.2, fill=LTGREY, line=BLUE, line_w=1.4)
add_text(slide, "EC2  t2.micro  (Ubuntu 22.04)  —  Docker Compose", 3.82, 1.55, 6.1, 0.38, size=10, bold=True, color=BLUE)

# nginx
add_rect(slide, 3.85, 2.0, 2.7, 0.9, fill=WHITE, line=INDIGO, line_w=0.8)
add_text(slide, "nginx:alpine",         3.98, 2.05, 2.4, 0.35, size=12, bold=True, color=NAVY)
add_text(slide, "Port 80 · Serves HTML\nProxies /api → backend:5000", 3.98, 2.4, 2.5, 0.42, size=8, color=GREY)

add_text(slide, "↓  proxy", 5.05, 2.97, 1.0, 0.28, size=9, italic=True, color=GREY)

# Flask
add_rect(slide, 3.85, 3.28, 2.7, 1.05, fill=WHITE, line=INDIGO, line_w=0.8)
add_text(slide, "Flask / Gunicorn",     3.98, 3.34, 2.4, 0.35, size=12, bold=True, color=NAVY)
add_text(slide, "Port 5000 · REST API\nOCR · Auth · Dashboard\n/upload /documents /dashboard", 3.98, 3.68, 2.5, 0.55, size=8, color=GREY)

add_text(slide, "↓  SQL", 5.05, 4.4, 1.0, 0.28, size=9, italic=True, color=GREY)

# Postgres
add_rect(slide, 3.85, 4.72, 2.7, 0.88, fill=WHITE, line=INDIGO, line_w=0.8)
add_text(slide, "PostgreSQL 15",        3.98, 4.78, 2.4, 0.35, size=12, bold=True, color=NAVY)
add_text(slide, "Port 5432 · Docker volume\nusers / documents / extracted_data", 3.98, 5.12, 2.5, 0.38, size=8, color=GREY)

# finsight-net brace
add_rect(slide, 7.35, 2.0, 2.5, 3.6, fill=LTBLUE, line=INDIGO, line_w=0.5)
add_text(slide, "finsight-net\n(Docker bridge\nnetwork)\n\nInternal only —\nnot reachable\nfrom internet", 7.45, 2.08, 2.3, 3.4, size=9, color=INDIGO)

# S3
add_rect(slide, 10.45, 2.1, 2.5, 1.1, fill=LTBLUE, line=INDIGO, line_w=1.0)
add_text(slide, "☁  Amazon S3",         10.58, 2.16, 2.25, 0.38, size=12, bold=True, color=NAVY)
add_text(slide, "File storage\nAES-256 encrypted\nPresigned preview URLs", 10.58, 2.52, 2.25, 0.55, size=8, color=BLUE)

# IAM
add_rect(slide, 10.45, 3.4, 2.5, 0.95, fill=LTBLUE, line=INDIGO, line_w=1.0)
add_text(slide, "🔑  IAM LabRole",       10.58, 3.46, 2.25, 0.38, size=12, bold=True, color=NAVY)
add_text(slide, "Instance Profile\nNo hardcoded keys", 10.58, 3.82, 2.25, 0.42, size=8, color=BLUE)

# Security Group
add_rect(slide, 10.45, 4.55, 2.5, 0.95, fill=LTBLUE, line=INDIGO, line_w=1.0)
add_text(slide, "🛡  Security Group",    10.58, 4.61, 2.25, 0.38, size=12, bold=True, color=NAVY)
add_text(slide, "Port 80 + 22 open\nPort 5000 blocked", 10.58, 4.97, 2.25, 0.42, size=8, color=BLUE)

# boto3 arrow
add_text(slide, "← boto3 →\nupload / download", 6.68, 3.48, 3.55, 0.55, size=8, italic=True, color=GREY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — AWS Services Deep Dive
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide)
header(slide, "AWS Services — What We Used & Why", "Every cloud service in the project explained")
foot(slide)

services = [
    ("EC2  t2.micro", "Virtual Server",
     "Runs all three Docker containers (nginx, Flask, PostgreSQL) on Ubuntu 22.04. Chosen over Lambda because Tesseract OCR needs OS-level binaries and ~500 MB of language data installed.",
     GREEN),
    ("Amazon S3", "Object Storage",
     "Every uploaded document is stored in S3 (not on EC2 disk). Files are AES-256 encrypted at rest. Frontend never accesses S3 directly — Flask generates a presigned URL (expires in 1 hour) for previews.",
     GREEN),
    ("IAM  LabRole", "Identity & Access",
     "The EC2 instance has a LabRole attached. boto3 fetches temporary credentials from the Instance Metadata Service automatically — no AWS keys are ever written in the code.",
     GREEN),
    ("Security Groups", "Firewall",
     "Port 80 open to all (web app). Port 22 open for SSH. Port 5000 (Flask) is NOT open — only nginx can reach it via the internal Docker network. This prevents direct API access.",
     GREEN),
    ("VPC", "Private Network",
     "All resources run inside the default VPC. Docker containers communicate over a private bridge network (finsight-net) that is not reachable from the internet.",
     GREEN),
    ("RDS PostgreSQL", "Managed Database",
     "Planned but replaced. AWS Academy deletes RDS instances on session expiry and blocks creating DB Subnet Groups (needs 2 AZs). Replaced with postgres:15-alpine in Docker with a named volume.",
     ORANGE),
    ("Lambda", "Serverless Trigger",
     "Designed as an async OCR trigger: S3 PUT event → Lambda → POST /process on Flask. Not deployed — Academy blocked the required execution role. OCR runs synchronously in the upload endpoint instead.",
     ORANGE),
]

row_h = 0.74
yy = 1.52
for svc, cat, desc, status_color in services:
    add_rect(slide, 0.35, yy, 12.65, row_h, fill=LTGREY, line=RGBColor(0xc5,0xd3,0xf0), line_w=0.3)
    add_rect(slide, 0.35, yy, 0.08, row_h, fill=status_color)
    add_text(slide, svc,  0.55, yy + 0.06, 2.1,  row_h - 0.12, size=10, bold=True, color=NAVY)
    add_text(slide, cat,  2.7,  yy + 0.06, 1.5,  row_h - 0.12, size=9,  italic=True, color=INDIGO)
    add_text(slide, desc, 4.35, yy + 0.06, 8.55, row_h - 0.12, size=9,  color=NAVY)
    yy += row_h + 0.02

add_text(slide, "● Active & working   ● Replaced / not deployed in Academy sandbox",
         0.35, 6.88, 10, 0.28, size=8, italic=True, color=GREY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK)
bg(slide, NAVY)
add_rect(slide, 0, 0, 0.2, 7.5, fill=INDIGO)

add_text(slide, "Thank You", 0.6, 1.3, 12, 1.2,
         size=64, bold=True, color=WHITE)
add_rect(slide, 0.6, 2.72, 4.2, 0.07, fill=INDIGO)

add_text(slide, "FinSight — Financial Document Intelligence Platform",
         0.6, 2.9, 11.5, 0.6, size=20, color=LTBLUE)

# summary pills
pills = [
    ("Upload", "PDF, PNG, JPG\nmulti-file"),
    ("OCR", "Tesseract\n300 DPI"),
    ("Extract", "Amounts, Dates\nVendors"),
    ("Categorise", "6 document\ntypes"),
    ("Dashboard", "Charts, trends\nmonth filter"),
    ("Cloud", "EC2, S3\nIAM, VPC"),
]
xx = 0.6
for label, desc in pills:
    add_rect(slide, xx, 3.75, 1.95, 1.4, fill=BLUE, line=INDIGO, line_w=0.8)
    add_text(slide, label, xx+0.1, 3.82, 1.72, 0.42, size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, xx, 4.24, 1.95, 0.04, fill=INDIGO)
    add_text(slide, desc,  xx+0.1, 4.32, 1.72, 0.7,  size=9, color=LTBLUE, align=PP_ALIGN.CENTER)
    xx += 2.05

add_text(slide, "Questions?", 0.6, 5.5, 6, 0.65, size=26, bold=True, color=INDIGO)
add_text(slide, "Built on AWS Academy  ·  EC2  ·  S3  ·  IAM  ·  VPC  ·  Docker",
         0.6, 6.25, 11, 0.42, size=13, italic=True, color=GREY)


# ── Save ───────────────────────────────────────────────────────────────────────
OUT = "FinSight_Presentation_v2.pptx"
prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
